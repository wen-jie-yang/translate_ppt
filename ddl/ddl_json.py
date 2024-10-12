import json
from pptx import Presentation
import openai

# 输入和输出文件路径
input_ppt = "/home/ywj/project/ppt/2.PPT/2.pptx"
output_ppt = "out3.pptx"

# 加载 PPT 文件
prs = Presentation(input_ppt)

# 使用 OpenAI 客户端进行翻译
client = openai.OpenAI(
    api_key="noneed",
    base_url="http://10.10.109.214:8800/v1"
)

def translate_text(text, target_language='zh'):
    """使用大模型翻译文本"""
    if not text.strip():
        return text  # 空文本直接返回

    response = client.chat.completions.create(
        temperature=0.2,
        model="glm-4-9b",
        messages=[
            {"role": "system", "content": f"你是一个专业的翻译专家。用户希望你将输入的文本翻译为{target_language}。"},
            {"role": "system", "content": "请准确翻译 JSON 结构中的文本，不要改变结构。只输出标准json格式的字符串，不输出任何其他内容。"},
            {"role": "user", "content": text}
        ]
    )
    # 获取翻译结果
    translated_text = response.choices[0].message.content
    return translated_text

# 将页面的 JSON 数据转为一个字符串，方便翻译
def create_json_string(page_data):
    json_string = json.dumps(page_data, ensure_ascii=False, indent=2)
    return json_string

# 遍历每一页，逐页翻译
for slide_index, slide in enumerate(prs.slides):
    page_data = {
        "slide_index": slide_index,
        "text_boxes": []
    }

    # 提取每个 slide 的 runs 文本并保存为 JSON 格式
    for shape_index, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            runs_data = []
            for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                for run_index, run in enumerate(paragraph.runs):
                    runs_data.append({
                        "run_index": run_index,
                        "paragraph_index": paragraph_index,
                        "text": run.text  # 记录每个 run 的文本
                    })
            # 将所有 runs 数据保存为 JSON 结构
            text_box_data = {
                "shape_index": shape_index,
                "runs": runs_data
            }
            page_data["text_boxes"].append(text_box_data)

    # 将该页的 JSON 数据转为字符串
    json_string = create_json_string(page_data)

    # 翻译该页的 JSON 数据
    translated_json_string = translate_text(json_string)

    # 将翻译后的 JSON 转回字典格式
    translated_page_data = json.loads(translated_json_string)
    print(translated_page_data)
    # 将翻译后的文本写回到该页
    for text_box_data in translated_page_data["text_boxes"]:
        shape = slide.shapes[text_box_data["shape_index"]]
        if shape.has_text_frame:
            # 遍历每个 run，并将翻译后的文本写回
            for run_data in text_box_data["runs"]:
                paragraph = shape.text_frame.paragraphs[run_data["paragraph_index"]]
                run = paragraph.runs[run_data["run_index"]]
                print(run_data["text"])
                run.text = run_data["text"]  # 替换为翻译后的文本

# 保存翻译后的 PPT
prs.save(output_ppt)
print("翻译完成！")
