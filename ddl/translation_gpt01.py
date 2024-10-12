import openai
from pptx import Presentation
from openai import OpenAI
from zhipuai import ZhipuAI

# 设置 OpenAI API 密钥
#

client = OpenAI(
    api_key="noneed",
    # base_url="https://open.bigmodel.cn/api/paas/v4/"
    base_url="http://10.10.109.214:8800/v1"

    # base_url="http://10.10.109.2:8800/v1"
)


def translate_text(text, target_language='zh'):
    """使用 OpenAI 的 API 翻译文本"""
    # 检查文本是否为空或仅包含空白字符
    if not text.strip():
        return text

    # 调用API进行翻译
    response = client.chat.completions.create(
        model="glm-4-9b",
        # model = "gpt-4o",
        temperature=0.2,
        messages=[
            {"role": "system",
             "content": f"你是一个计算机相关行业的专业翻译软件，你需要将用户发送给你文本翻译为 {target_language} 语言。注意判断：专业名词和数学公式不需要翻译。注意你的翻译过程是交互式的，用户发一段文本，你就翻译一次。"},
            {"role": "system",
             "content": f"请不要输出多余的话语！只输出翻译结果,原文中含有的'<split_sign>'字符串代表分割符，是用户做的标识，译文中你需要在适当的位置保留同等数量的'<split_sign>'字符串。"},

            {"role": "user", "content": f"{text}"}
        ]
    )
    # 假设API返回的翻译结果存储在response.choices[0].message.content中
    return response.choices[0].message.content


def translate_ppt(input_ppt, output_ppt):
    """每页处理完文本后立即翻译并放回各自的 run 中"""
    prs = Presentation(input_ppt)

    for slide in prs.slides:
        # 存储当前页的所有文本
        all_runs_text = []
        slide_run_indices = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run_index, run in enumerate(paragraph.runs):
                        original_text = run.text
                        if original_text.strip():
                            all_runs_text.append(original_text)
                            slide_run_indices.append((paragraph, run_index))  # 记录对应的 paragraph 和 run 的索引

        # 只有在有内容时才进行翻译处理
        if all_runs_text:
            # 合并所有 run 的文本，用 `\\\\` 作为分隔符
            combined_text = '<split_sign>'.join(all_runs_text)
            print(combined_text)
            # 调用翻译函数进行翻译
            translated_text = translate_text(combined_text)
            print(translated_text)
            # 拆分翻译后的文本
            translated_runs = translated_text.split('<split_sign>')

            # 将翻译后的文本放回各自的 run 中
            for (paragraph, run_index), translated_run_text in zip(slide_run_indices, translated_runs):
                paragraph.runs[run_index].text = translated_run_text

    # 保存翻译后的 PPT 文件
    prs.save(output_ppt)


# 输入和输出 PPT 文件路径
input_ppt = r'test.pptx'  # 替换为你的输入文件路径
output_ppt = r'out.pptx'  # 替换为你的输出文件路径

translate_ppt(input_ppt, output_ppt)
print("翻译完成！")
