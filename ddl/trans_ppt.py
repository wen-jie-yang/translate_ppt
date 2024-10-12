import gradio as gr
import tempfile
import os
from pptx import Presentation
import openai

# 配置 OpenAI 客户端
client = openai.OpenAI(
    api_key="noneed",  # 替换为你的实际 API 密钥
    base_url="http://10.10.109.214:8800/v1"  # 替换为你的 API 服务地址
)

def translate_text(text, target_language='zh'):
    """使用大模型翻译文本，并保留 <split_sign> 分隔符"""
    try:
        if not text.strip():
            return text  # 空文本直接返回

        response = client.chat.completions.create(
            model="glm-4-9b",
            temperature=0.2,
            messages=[
                {"role": "system", "content": f"你是一个专业的翻译专家。用户希望你将输入的文本翻译为{target_language}。"},
                {"role": "system", "content": "注意，文本中出现的 '<split_sign>' 是自定义的分隔符，表示文本段落的划分。"},
                {"role": "system", "content": "请确保所有 '<split_sign>' 保持原样，不要翻译或删除这些分隔符。"},
                {"role": "system", "content": "请确保所有 '<split_sign>' 保持原样，不要翻译、删除或修改任何 '<split_sign>'，并且它们是文本段落的标识符。"},
                {"role": "system", "content": "输出时只需给出翻译结果，并保留所有 '<split_sign>'，不需要输出任何解释性内容。"},
                {"role": "user", "content": text}
            ]
        )
        translated_text = response.choices[0].message.content
        return translated_text
    except Exception as e:
        print(f"翻译时发生错误: {e}")
        return text  # 如果翻译失败，返回原文

def process_ppt(input_ppt):
    """处理PPT文件并进行翻译"""
    try:
        # 确保文件是 .pptx 格式
        if not input_ppt.name.endswith('.pptx'):
            raise ValueError("仅支持 PPTX 文件格式")

        # 加载 PPT 文件
        prs = Presentation(input_ppt)

        # 提取每页的所有 runs 文本内容
        pages = []
        for slide_index, slide in enumerate(prs.slides):
            page_data = {
                "slide_index": slide_index,
                "text_boxes": []
            }
            for shape_index, shape in enumerate(slide.shapes):
                if shape.has_text_frame:
                    runs_text = []
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            runs_text.append(run.text)
                    # 将每个 shape 的所有 run.text 拼接为一个字符串
                    text_box_data = {
                        "shape_index": shape_index,
                        "runs_text": runs_text  # 保留原始的 runs 列表
                    }
                    page_data["text_boxes"].append(text_box_data)
            pages.append(page_data)

        # 拼接每页所有文本框中的 run.text，用 <split_sign> 作为分隔符
        page_texts = []
        for page_data in pages:
            all_runs = []
            for text_box in page_data["text_boxes"]:
                all_runs.extend(text_box["runs_text"])
            # 将所有 run.text 拼接为一个字符串
            combined_text = '<split_sign>'.join(all_runs)
            page_texts.append(combined_text)

        # 翻译每页的文本
        translated_pages = []
        for combined_text in page_texts:
            translated_combined_text = translate_text(combined_text)
            # 将翻译结果按 <split_sign> 分割
            translated_runs = translated_combined_text.split('<split_sign>')
            translated_pages.append(translated_runs)

        # 将翻译结果写回到原来的位置
        for page_data, translated_runs in zip(pages, translated_pages):
            slide = prs.slides[page_data["slide_index"]]
            run_index = 0
            for text_box_data in page_data["text_boxes"]:
                shape = slide.shapes[text_box_data["shape_index"]]
                if shape.has_text_frame:
                    # 更新每个 run 的文本
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run_index < len(translated_runs):
                                run.text = translated_runs[run_index]
                                run_index += 1

        # 使用当前工作目录下的自定义目录保存输出文件
        output_dir = "output_files"
        os.makedirs(output_dir, exist_ok=True)  # 如果不存在则创建
        output_path = os.path.join(output_dir, "translated_output.pptx")  # 输出文件路径
        prs.save(output_path)
        return output_path
    except Exception as e:
        print(f"处理PPT时出错: {e}")
        raise e  # 如果有错误，抛出异常

def translate_ppt(input_ppt):
    """Gradio 接口：处理输入的PPT文件并输出翻译后的PPT文件"""
    try:
        if input_ppt is None:
            return "未收到有效的 PPT 文件"

        # 处理并返回翻译后的 PPT 文件路径
        translated_ppt_path = process_ppt(input_ppt)
        return translated_ppt_path
    except Exception as e:
        print(f"翻译过程中出现错误: {e}")
        return f"翻译过程中出现错误: {e}"

# 创建自定义临时目录来替代 /tmp
custom_temp_dir = os.path.join(os.getcwd(), "custom_temp")
os.makedirs(custom_temp_dir, exist_ok=True)

# 设置全局的临时目录为 custom_temp_dir
tempfile.tempdir = custom_temp_dir

# 使用 Gradio 创建界面
iface = gr.Interface(
    fn=translate_ppt,  # 翻译函数
    inputs=gr.File(label="上传PPT文件"),  # 输入为PPT文件
    outputs=gr.File(label="下载翻译后的PPT文件"),  # 输出为PPT文件
    title="PPT 翻译器",
    description="上传一个PPT文件，自动翻译为中文并返回翻译后的PPT文件"
)

# 启动 Gradio 服务，指定端口 18901
iface.launch(server_name="0.0.0.0", server_port=18901)
