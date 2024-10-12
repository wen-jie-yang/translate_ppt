from pptx import Presentation
import openai

# 输入和输出文件路径
input_ppt = "/home/ywj/project/ppt/ddl/m3-Image_Content_Analysis.pptx"
output_ppt = "out22.pptx"

# 配置 OpenAI 客户端
client = openai.OpenAI(
    api_key="noneed",
    base_url="http://10.10.109.214:8800/v1"
)

def translate_text(text, target_language='zh'):
    """使用大模型翻译文本，并保留 <split_sign> 分隔符"""
    if not text.strip():
        return text  # 空文本直接返回

    response = client.chat.completions.create(
        model="glm-4-9b",
        temperature=0.2,
        messages=[
            {"role": "system", "content": f"你是一个专业的翻译专家。用户希望你将输入的文本翻译为{target_language}。"},
            {"role": "system", "content": "注意，文本中出现的 '<split_sign>' 是自定义的分隔符，表示文本段落的划分。"},
            {"role": "system", "content": "请确保所有 '<split_sign>' 保持原样，不要翻译或删除这些分隔符。"},
            {"role": "system", "content": "请确保所有 '<split_sign>' 保持原样，不要翻译、删除或修改任何 '<split_sign>'，并且它们是文本段落的标识符。"},
            {"role": "system", "content": "输出时只需给出翻译结果，并保留所有 '<split_sign>'，不需要输出任何解释性内容。"},
            {"role": "user", "content": text}
        ]
    )
    # 获取翻译结果
    translated_text = response.choices[0].message.content
    return translated_text

def process_ppt(input, output):
    # 加载 PPT 文件
    prs = Presentation(input)

    # 提取每页的所有 runs 文本内容
    pages = []
    for slide_index, slide in enumerate(prs.slides):
        page_data = {
            "slide_index": slide_index,
            "text_boxes": []
        }
        for shape_index, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                runs_text = []
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        runs_text.append(run.text)
                # 将每个 shape 的所有 run.text 拼接为一个字符串
                text_box_data = {
                    "shape_index": shape_index,
                    "runs_text": runs_text  # 保留原始的 runs 列表
                }
                page_data["text_boxes"].append(text_box_data)
        pages.append(page_data)

    # 拼接每页所有文本框中的 run.text，用 <split_sign> 作为分隔符
    page_texts = []
    for page_data in pages:
        all_runs = []
        for text_box in page_data["text_boxes"]:
            all_runs.extend(text_box["runs_text"])
        # 将所有 run.text 拼接为一个字符串
        combined_text = '<split_sign>'.join(all_runs)
        page_texts.append(combined_text)

    # 翻译每页的文本
    translated_pages = []
    for combined_text in page_texts:
        translated_combined_text = translate_text(combined_text)
        # 将翻译结果按 <split_sign> 分割
        translated_runs = translated_combined_text.split('<split_sign>')
        translated_pages.append(translated_runs)

    # 将翻译结果写回到原来的位置
    for page_data, translated_runs in zip(pages, translated_pages):
        slide = prs.slides[page_data["slide_index"]]
        run_index = 0
        for text_box_data in page_data["text_boxes"]:
            shape = slide.shapes[text_box_data["shape_index"]]
            if shape.has_text_frame:
                # 更新每个 run 的文本
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run_index < len(translated_runs):
                            run.text = translated_runs[run_index]
                            run_index += 1
    # 保存翻译后的 PPT
    prs.save(output)
    print("翻译完成！")
    return output


