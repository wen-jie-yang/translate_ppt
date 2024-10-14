import json
from pptx import Presentation
import openai

# 输入和输出文件路径
input_ppt = "/home/ywj/project/ppt/ddl/m3-Image_Content_Analysis.pptx"  # 替换为实际路径
json_output = "output.json"
output_ppt = "output_translated.pptx"

# 使用 OpenAI 客户端进行翻译
client = openai.OpenAI(
    api_key="noneed",  # 替换为你的实际 API 密钥
    base_url="http://10.10.109.214:8800/v1"  # 替换为你的 API 服务地址
)

def extract_ppt_to_json(input_ppt, json_output):
    """提取 PPT 内容到 JSON 文件，精确到每一个 run"""
    prs = Presentation(input_ppt)
    presentation_data = []

    for slide_index, slide in enumerate(prs.slides):
        page_data = {
            "slide_index": slide_index,
            "text_boxes": []
        }

        for shape_index, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                runs_data = []
                for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                    for run_index, run in enumerate(paragraph.runs):
                        runs_data.append({
                            "run_index": run_index,
                            "paragraph_index": paragraph_index,
                            "text": run.text  # 精确到每个 run
                        })
                text_box_data = {
                    "shape_index": shape_index,
                    "runs": runs_data
                }
                page_data["text_boxes"].append(text_box_data)
        
        presentation_data.append(page_data)

    # 保存提取的 JSON 数据
    with open(json_output, 'w', encoding='utf-8') as json_file:
        json.dump(presentation_data, json_file, ensure_ascii=False, indent=2)
    print(f"PPT 提取到 JSON 文件成功，路径：{json_output}")

def translate_json_page(page_data, target_language='zh'):
    """使用大模型翻译单个页面的 JSON 数据"""
    json_string = json.dumps(page_data, ensure_ascii=False, indent=2)

    prompt = f"""
    你是一个专业的翻译专家。用户希望你将输入的 JSON 数据中的文本翻译为 {target_language}，并保持原始的 JSON 结构。请你将以下 JSON 结构中的所有文本进行翻译，保持 JSON 结构完全不变，只翻译文本字段：
    {json_string}
    """

    response = client.chat.completions.create(
        temperature=0.2,
        model="glm-4-9b",
        messages=[
            {"role": "system", "content": "你是一个JSON数据翻译专家。"},
            {"role": "user", "content": prompt}
        ]
    )
    translated_json = response.choices[0].message.content
    print("大模型返回的内容：", translated_json)
        # 移除可能的代码块标记（如 ```json 或 ```）
    if translated_json.startswith("```json"):
        translated_json = translated_json[7:].strip()
    if translated_json.endswith("```"):
        translated_json = translated_json[:-3].strip()
    return json.loads(translated_json)  # 返回翻译后的 JSON 数据

def translate_json_file(json_input, json_output_translated):
    """按页翻译 JSON 文件的内容，并保存翻译后的 JSON 文件"""
    with open(json_input, 'r', encoding='utf-8') as json_file:
        presentation_data = json.load(json_file)

    translated_presentation_data = []

    for page_data in presentation_data:
        # 将每个页面的 JSON 数据直接传递给大模型
        translated_page_data = translate_json_page(page_data)
        translated_presentation_data.append(translated_page_data)

    # 保存翻译后的 JSON 数据
    with open(json_output_translated, 'w', encoding='utf-8') as json_file:
        json.dump(translated_presentation_data, json_file, ensure_ascii=False, indent=2)
    print(f"翻译后的 JSON 文件已保存，路径：{json_output_translated}")

def replace_ppt_with_translated_json(json_input, output_ppt, original_ppt):
    """根据翻译后的 JSON 文件替换 PPT 内容"""
    prs = Presentation(original_ppt)

    with open(json_input, 'r', encoding='utf-8') as json_file:
        presentation_data = json.load(json_file)

    for page_data in presentation_data:
        slide = prs.slides[page_data["slide_index"]]
        
        for text_box_data in page_data["text_boxes"]:
            shape = slide.shapes[text_box_data["shape_index"]]
            if shape.has_text_frame:
                for run_data in text_box_data["runs"]:
                    paragraph = shape.text_frame.paragraphs[run_data["paragraph_index"]]
                    run = paragraph.runs[run_data["run_index"]]
                    run.text = run_data["text"]  # 替换为翻译后的文本

    # 保存翻译后的 PPT
    prs.save(output_ppt)
    print(f"翻译后的 PPT 已保存，路径：{output_ppt}")

# Step 1: 提取 PPT 内容到 JSON
extract_ppt_to_json(input_ppt, json_output)

# Step 2: 翻译 JSON 文件内容并保存
translated_json_output = "translated_output.json"
translate_json_file(json_output, translated_json_output)

# Step 3: 将翻译后的 JSON 内容替换回 PPT
replace_ppt_with_translated_json(translated_json_output, output_ppt, input_ppt)
