import openai
from pptx import Presentation
from openai import OpenAI
import zhipuai
# 设置 OpenAI API 密钥
# import time
 
# # 记录开始时间
# start_time = time.time()

client = OpenAI(
    api_key="noneed",
    # base_url="https://open.bigmodel.cn/api/paas/v4/"
    base_url="http://10.10.109.214:8800/v1"
) 


def translate_text(text, target_language='zh'):
   """使用 OpenAI 的 API 翻译文本"""
   # 检查文本是否为空或仅包含空白字符
   if not text.strip():
       return text

   # 以下是调用API的伪代码，因为没有具体的API实现
   response = client.chat.completions.create(
    #    model="glm-4-0520",
       model="glm-4-9b",

       messages=[
           {"role":"system","content":f"你是一个计算机相关行业的专业翻译软件，你需要将用户发送给你文本翻译为 {target_language} 语言。注意你的翻译过程是交互式的，用户发一段文本，你就翻译一次。由于你是专业翻译软件，在遇到专业名词等的时候，你需要自主判断是否应该翻译。注意，由于你是一个软件，所以你不会与用户交流，不会输出任何提示性词汇，你只会输出翻译后的文本,如果无法翻译，直接输出原文。"},
           {"role": "user", "content": f"{text}"}
       ]
   )
   # 假设API返回的翻译结果存储在response.choices[0].message.content中
   return response.choices[0].message.content

def translate_ppt(input_ppt, output_ppt):
   """翻译 PPT 文件中的所有文本"""
   prs = Presentation(input_ppt)
   for slide in prs.slides:
       for shape in slide.shapes:
           if shape.has_text_frame:
               text_frame = shape.text_frame
               for paragraph in text_frame.paragraphs:
                   for run in paragraph.runs:
                       original_text = run.text
                       # 仅翻译非空字符串
                       if original_text.strip():
                           translated_text = translate_text(original_text)
                           run.text = translated_text

   prs.save(output_ppt)

# 输入和输出 PPT 文件路径
input_ppt = r'm3-Image_Content_Analysis.pptx'  # 替换为你的输入文件路径
output_ppt = r'output.pptx'  # 替换为你的输出文件路径

translate_ppt(input_ppt, output_ppt)
print("翻译完成！")
# 你的代码逻辑
# ... 
# 记录结束时间
# end_time = time.time() 

# # 计算并打印运行时长
# print(f"运行时间：{end_time - start_time} 秒")

