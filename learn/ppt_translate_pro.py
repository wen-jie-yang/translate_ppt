import openai
from pptx import Presentation
import json

# 设置 OpenAI 客户端
client = openai.OpenAI(
    api_key="noneed",
    base_url="http://10.10.109.214:8800/v1"
)

def translate_text(text, target_language='zh'):
    """使用 OpenAI 的 API 翻译整段文本"""
    if not text.strip():
        return text  # 空文本直接返回

    # 调用大模型进行翻译
    response = client.chat.completions.create(
        model="glm-4-9b",
        messages=[
            {"role": "system", "content": f"你是一个专业的翻译软件，将用户的文本翻译成 {target_language}。"},
            {"role": "user", "content": text}
        ]
    )
    return response.choices[0].message.content

def extract_text_from_ppt(prs):
    """提取每页的文本框内容并存储为 JSON 格式"""
    pages = []
    for slide_index, slide in enumerate(prs.slides):
        page_data = {
            "slide_index": slide_index,
            "text_boxes": []
        }
        for shape_index, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                text_content = []
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_content.append(run.text)
                # 将文本内容和位置存储
                page_data["text_boxes"].append({
                    "shape_index": shape_index,
                    "original_text": "\n".join(text_content)
                })
        pages.append(page_data)
    return pages

def insert_translated_text(prs, pages, translations):
    """将翻译后的文本放回到原来的位置"""
    for page_data, translation in zip(pages, translations):
        slide = prs.slides[page_data["slide_index"]]
        for text_box, translated_text in zip(page_data["text_boxes"], translation["translated_texts"]):
            shape = slide.shapes[text_box["shape_index"]]
            if shape.has_text_frame:
                # 更新文本框的内容
                shape.text = translated_text

def translate_ppt(input_ppt, output_ppt):
    """翻译 PPT 文件中的所有文本，并替换翻译结果"""
    prs = Presentation(input_ppt)

    # 提取每页的文本并存储为 JSON
    pages = extract_text_from_ppt(prs)

    # 拼接每页的所有文本内容，进行一次性翻译
    translations = []
    for page_data in pages:
        all_text = "\n\n".join([text_box["original_text"] for text_box in page_data["text_boxes"]])
        translated_text = translate_text(all_text)
        # 分割翻译后的文本，根据原有的分段数量
        translated_texts = translated_text.split("\n\n")
        translations.append({
            "slide_index": page_data["slide_index"],
            "translated_texts": translated_texts
        })

    # 将翻译后的文本放回到原来的位置
    insert_translated_text(prs, pages, translations)

    # 保存翻译后的 PPT
    prs.save(output_ppt)
    print("翻译完成！")

# 输入和输出 PPT 文件路径
input_ppt = r'm3-Image_Content_Analysis.pptx'  # 替换为你的输入文件路径
output_ppt = r'output.pptx'  # 替换为你的输出文件路径

translate_ppt(input_ppt, output_ppt)
