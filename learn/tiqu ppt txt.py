import json # 导入 Python 自带的 json 模块，用于处理 JSON 格式数据
from pptx import Presentation # 从 pptx 库中导入 Presentation 类，用于操作 PPT 文件

def extract_text_to_json(ppt_path, json_path):
    # 打开指定路径的 PPT 文件
    prs = Presentation(ppt_path)

    # 创建一个空列表，用于存储所有幻灯片的文本数据   
    ppt_data = []
    
    # 遍历每一页幻灯片，enumerate 用于同时获取索引（slide_index）和幻灯片对象（slide）
    for slide_index, slide in enumerate(prs.slides):
        # 创建一个字典，存储当前幻灯片的编号和所有文本框的信息
        slide_info = {"slide_number": slide_index + 1, "text_boxes": []}

        # 遍历当前幻灯片中的每个形状（可能是文本框、图片等）
        for shape_index, shape in enumerate(slide.shapes):
            # 检查该形状是否有文本框
            if shape.has_text_frame:
                text_content = ""  # 用于存储该文本框的所有文本内容

                # 遍历文本框中的每个段落
                for paragraph in shape.text_frame.paragraphs:
                    
                    # 将段落的文本添加到 text_content 中，使用 '\n' 表示换行
                    text_content += paragraph.text + "\n"
                
                # 创建一个字典，存储该文本框的 ID（编号）和文本内容
                text_box_info = {
                    "box_id": shape_index + 1,  # ID 为当前形状在幻灯片中的索引 +1
                    "text": text_content.strip()  # strip() 方法去掉首尾的空白字符
                }
                
                # 将该文本框信息添加到当前幻灯片的信息中
                slide_info["text_boxes"].append(text_box_info)

        # 将当前幻灯片的信息添加到整个 PPT 数据列表中
        ppt_data.append(slide_info)

    # 打开指定路径的 JSON 文件，并将 PPT 数据写入其中
    with open(json_path, "w", encoding="utf-8") as ison_file:
        # json.dump() 将 Python 数据对象（ppt_data）写入到 JSON 文件中
        # indent=4 使 JSON 数据更易读，ensure_ascii=False 保持中文字符的正确显示
        json.dump(ppt_data, ison_file, indent=4, ensure_ascii=False)
    
    # 提示用户操作完成
    print(f"Text data extracted to {json_path}")

# 使用示例，指定要读取的 PPT 文件路径和保存 JSON 文件路径
ppt_path = r"D:\learn\learn_JSON\m3-Image_Content_Analysis.pptx"  # 替换为你的输入文件路径
json_path = r"D:\learn\learn_JSON\ppt_text_data.json"  # 替换为输出 JSON 文件路径

# 调用函数，执行文本提取操作
extract_text_to_json(ppt_path, json_path)
